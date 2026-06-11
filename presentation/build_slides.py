from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

os.makedirs("presentation", exist_ok=True)

# ── Colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0F, 0x1A)   # near-black navy
INDIGO   = RGBColor(0x63, 0x66, 0xF1)   # primary accent
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0xA0, 0xA8, 0xC0)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
PURPLE   = RGBColor(0xA7, 0x8B, 0xFA)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, x, y, w, h,
          size=20, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def multiline(slide, lines, x, y, w, h,
              size=16, color=WHITE, spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, sz, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(sz)
        run.font.bold  = bold
        run.font.color.rgb = col if col else color
    return tb

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════
s1 = add_slide(); bg(s1)

# Indigo accent bar left
rect(s1, 0, 0, 0.06, 7.5, INDIGO)

# Title
txbox(s1, "TridentHook", 0.4, 1.8, 8, 1.4, size=60, bold=True, color=WHITE)

# Subtitle
txbox(s1, "Dynamic Fee Protection for Uniswap v4 Liquidity Providers",
      0.4, 3.3, 9, 0.8, size=24, color=GREY)

# Tag line
txbox(s1, "Research-driven  ·  Fully On-Chain  ·  Autonomous",
      0.4, 4.2, 9, 0.5, size=16, color=INDIGO)

# Built on badges
txbox(s1, "Built on Unichain Sepolia   |   Powered by Reactive Network   |   Uniswap v4 Hook",
      0.4, 6.6, 12, 0.5, size=13, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM (research backed)
# ═══════════════════════════════════════════════════════════════════════════
s2 = add_slide(); bg(s2)
rect(s2, 0, 0, 0.06, 7.5, YELLOW)

txbox(s2, "The Problem Is Proven — Not Assumed", 0.4, 0.25, 12, 0.8,
      size=32, bold=True, color=WHITE)
txbox(s2, "6 peer-reviewed papers. Same conclusion: LP capital is being systematically extracted.",
      0.4, 1.0, 12, 0.45, size=15, color=GREY)

# 4 quote cards  (2×2 grid)
cards = [
    ("\"Fees do not cover arbitrage losses in most of\n the largest Uniswap pools studied.\"",
     "Fritsch & Canidio — arXiv 2404.05803, 2024", "→ Layer 1: Arb Detector"),
    ("\"LVR is the dominant LP loss. Pools need 10%\n daily turnover at 30 bps just to break even.\"",
     "Milionis, Moallemi, Roughgarden & Zhang — arXiv 2208.06046, 2022", "→ Layer 1: Arb Detector"),
    ("\"LP positions are short gamma — risk spikes\n sharply at range boundaries.\"",
     "Impermanent Loss in Uniswap v3 — arXiv 2111.09192", "→ Layer 2: Range Guardian"),
    ("\"JIT bots dilute regular LP fees by 85% on average.\n $750B in JIT volume extracted per year.\"",
     "The Paradox of JIT Liquidity — arXiv 2311.18164, 2024", "→ Layer 3: IL Reserve Vault"),
]

positions = [(0.35, 1.65), (6.85, 1.65), (0.35, 4.35), (6.85, 4.35)]
for (quote, source, layer), (cx, cy) in zip(cards, positions):
    rect(s2, cx, cy, 6.1, 2.55, RGBColor(0x16, 0x1A, 0x2E))
    rect(s2, cx, cy, 0.06, 2.55, YELLOW)
    txbox(s2, quote,  cx+0.15, cy+0.12, 5.85, 1.15, size=13, color=WHITE)
    txbox(s2, source, cx+0.15, cy+1.35, 5.85, 0.45, size=10, color=GREY)
    txbox(s2, layer,  cx+0.15, cy+1.95, 5.85, 0.4,  size=11, bold=True, color=YELLOW)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════
s3 = add_slide(); bg(s3)
rect(s3, 0, 0, 0.06, 7.5, GREEN)

txbox(s3, "TridentHook: Three Layers of Protection", 0.4, 0.25, 12, 0.8,
      size=32, bold=True, color=WHITE)
txbox(s3, "Every mechanism is the direct on-chain implementation of what the research prescribes.",
      0.4, 0.95, 12, 0.4, size=14, color=GREY)

layers = [
    (INDIGO,  "Layer 1 — Arb Detector",
     "Reactive Network monitors Chainlink oracle live. When pool price diverges from fair value,\nfee spikes automatically — arbitrage becomes expensive, not free.",
     "Research: arXiv 2208.06046 · 2404.05803"),
    (PURPLE,  "Layer 2 — Range Guardian",
     "Detects when price approaches LP range boundaries. Fee elevates proportionally to gamma\nexposure — LPs earn the most exactly when they are at most risk.",
     "Research: arXiv 2111.09192 · 2407.05146"),
    (GREEN,   "Layer 3 — IL Reserve Vault",
     "A share of every swap fee flows into an on-chain reserve. LPs claim compensation on exit.\nLoyalty factor rewards long-term LPs over JIT bots.",
     "Research: arXiv 2410.00854 · 2502.04097 · 2311.18164"),
]

for i, (col, title, desc, research) in enumerate(layers):
    y = 1.6 + i * 1.75
    rect(s3, 0.35, y, 12.6, 1.55, RGBColor(0x16, 0x1A, 0x2E))
    rect(s3, 0.35, y, 0.06, 1.55, col)
    txbox(s3, title,    0.55, y+0.1,  5,    0.45, size=17, bold=True, color=col)
    txbox(s3, desc,     0.55, y+0.52, 8.5,  0.7,  size=13, color=WHITE)
    txbox(s3, research, 9.5,  y+0.1,  3.3,  0.45, size=10, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s4 = add_slide(); bg(s4)
rect(s4, 0, 0, 0.06, 7.5, PURPLE)

txbox(s4, "Fully Autonomous — No Bots, No Keepers", 0.4, 0.25, 12, 0.8,
      size=32, bold=True, color=WHITE)
txbox(s4, "Reactive Network replaces off-chain infrastructure with deterministic on-chain event execution.",
      0.4, 0.95, 12, 0.4, size=14, color=GREY)

# Left box — Unichain
rect(s4, 0.35, 1.55, 5.5, 5.5, RGBColor(0x16, 0x1A, 0x2E))
rect(s4, 0.35, 1.55, 5.5, 0.45, RGBColor(0x1E, 0x24, 0x40))
txbox(s4, "Unichain Sepolia", 0.5, 1.6, 5, 0.35, size=13, bold=True, color=INDIGO)

unichain_steps = [
    "① Swap event emitted by PoolManager",
    "② ReactiveAdapter receives callback",
    "③ TridentHook.primeDeviation() called",
    "④ Dynamic fee computed + applied",
    "⑤ Fee split: LPs + IL Reserve Vault",
]
for i, step in enumerate(unichain_steps):
    txbox(s4, step, 0.5, 2.15 + i*0.78, 5.1, 0.6, size=13, color=WHITE)

# Right box — Reactive Network
rect(s4, 7.48, 1.55, 5.5, 5.5, RGBColor(0x16, 0x1A, 0x2E))
rect(s4, 7.48, 1.55, 5.5, 0.45, RGBColor(0x1E, 0x24, 0x40))
txbox(s4, "Reactive Network (Lasna)", 7.6, 1.6, 5.2, 0.35, size=13, bold=True, color=PURPLE)

reactive_steps = [
    "① TridentReactive subscribes to events",
    "② Swap + AnswerUpdated detected",
    "③ Deviation vs Chainlink computed",
    "④ Gamma score from LP boundaries",
    "⑤ Callback emitted → Unichain",
]
for i, step in enumerate(reactive_steps):
    txbox(s4, step, 7.6, 2.15 + i*0.78, 5.2, 0.6, size=13, color=WHITE)

# Arrow between
txbox(s4, "⟵  On-chain callback  ⟶", 5.9, 3.8, 1.5, 0.8,
      size=11, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════
s5 = add_slide(); bg(s5)
rect(s5, 0, 0, 0.06, 7.5, INDIGO)

txbox(s5, "Live Demo", 0.4, 0.8, 8, 1.0, size=56, bold=True, color=WHITE)
txbox(s5, "Everything live on Unichain Sepolia · 146 tests passing",
      0.4, 1.85, 10, 0.45, size=15, color=GREY)

steps = [
    ("01", "Connect wallet on Unichain Sepolia"),
    ("02", "Add liquidity — position appears instantly"),
    ("03", "Trigger a swap — watch dynamic fee activate"),
    ("04", "Arb premium spikes when oracle gap detected"),
    ("05", "IL Reserve Vault accumulates with every trade"),
]
for i, (num, desc) in enumerate(steps):
    y = 2.55 + i * 0.88
    rect(s5, 0.35, y, 0.55, 0.62, INDIGO)
    txbox(s5, num,  0.37, y+0.08, 0.5,  0.45, size=15, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s5, desc, 1.05, y+0.1,  11,   0.45, size=17, color=WHITE)

txbox(s5, "Capture rate panel shows real-time state:  NORMAL  →  ELEVATED  →  MAX",
      0.35, 7.0, 12, 0.38, size=13, color=INDIGO)

# ── Save ────────────────────────────────────────────────────────────────────
out = "presentation/TridentHook_Demo.pptx"
prs.save(out)
print(f"Saved: {out}")
